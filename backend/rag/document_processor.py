"""
Document Processor — Multi-format document parsing
Mendukung: PDF, DOCX, XLSX, CSV, PPTX, TXT, MD, JSON
"""
import asyncio
import csv
import json
import os
from datetime import datetime
from pathlib import Path
from typing import List, Optional
import structlog

log = structlog.get_logger()

# Timeout per dokumen dalam detik
PROCESS_TIMEOUT = 45


class ProcessedChunk:
    """Representasi satu chunk teks dari dokumen."""
    def __init__(self, content: str, metadata: dict):
        self.page_content = content
        self.metadata = metadata


class DocumentProcessor:
    """
    Multi-format document processor.
    Digunakan oleh RAG engine untuk ekstrak teks dari berbagai format file.
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".csv", ".json", ".xlsx", ".xls", ".pptx", ".ppt"}

    @staticmethod
    def is_supported(file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        return ext in DocumentProcessor.SUPPORTED_EXTENSIONS

    @staticmethod
    def extract_metadata(file_path: str) -> dict:
        """Extract file metadata: name, size, date, extension."""
        p = Path(file_path)
        stat = p.stat() if p.exists() else None
        return {
            "filename": p.name,
            "extension": p.suffix.lower(),
            "file_size_bytes": stat.st_size if stat else 0,
            "file_size_kb": round(stat.st_size / 1024, 1) if stat else 0,
            "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat() if stat else None,
            "created_at": datetime.now().isoformat(),
        }

    # ── Synchronous extractors ──────────────────────────────────

    @staticmethod
    def _extract_pdf(file_path: str) -> List[str]:
        """Extract text from PDF, page by page."""
        pages = []
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            log.info("PDF loaded", pages=len(reader.pages), file=file_path)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    pages.append(f"[Halaman {i+1}]\n{text.strip()}")
        except Exception as e:
            log.error("PDF extract error", file=file_path, error=str(e))
            raise ValueError(f"Gagal membaca PDF: {e}")
        return pages

    @staticmethod
    def _extract_docx(file_path: str) -> List[str]:
        """Extract text from DOCX."""
        try:
            import docx
            doc = docx.Document(file_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            log.info("DOCX loaded", paragraphs=len(paragraphs), file=file_path)
            return paragraphs
        except Exception as e:
            log.error("DOCX extract error", file=file_path, error=str(e))
            raise ValueError(f"Gagal membaca DOCX: {e}")

    @staticmethod
    def _extract_xlsx(file_path: str) -> List[str]:
        """Extract text from XLSX/XLS spreadsheet."""
        rows = []
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append(" | ".join(cells))
            wb.close()
            log.info("XLSX loaded", rows=len(rows), file=file_path)
        except Exception as e:
            log.error("XLSX extract error", file=file_path, error=str(e))
            raise ValueError(f"Gagal membaca XLSX: {e}")
        return rows

    @staticmethod
    def _extract_csv(file_path: str) -> List[str]:
        """Extract text from CSV."""
        rows = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                for row in reader:
                    if any(cell.strip() for cell in row):
                        rows.append(" | ".join(row))
            log.info("CSV loaded", rows=len(rows), file=file_path)
        except Exception as e:
            log.error("CSV extract error", file=file_path, error=str(e))
            raise ValueError(f"Gagal membaca CSV: {e}")
        return rows

    @staticmethod
    def _extract_pptx(file_path: str) -> List[str]:
        """Extract text from PPTX/PPT slides."""
        slides = []
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))
            log.info("PPTX loaded", slides=len(slides), file=file_path)
        except Exception as e:
            log.error("PPTX extract error", file=file_path, error=str(e))
            raise ValueError(f"Gagal membaca PPTX: {e}")
        return slides

    @staticmethod
    def _extract_text(file_path: str) -> List[str]:
        """Extract text from TXT/MD files."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            # Split by paragraphs
            paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
            log.info("TXT/MD loaded", paragraphs=len(paragraphs), file=file_path)
            return paragraphs
        except Exception as e:
            log.error("TXT extract error", file=file_path, error=str(e))
            raise ValueError(f"Gagal membaca TXT: {e}")

    @staticmethod
    def _extract_json(file_path: str) -> List[str]:
        """Extract text from JSON files."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                data = json.load(f)
            text = json.dumps(data, ensure_ascii=False, indent=2)
            # Split into chunks of ~1000 chars
            chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
            log.info("JSON loaded", chunks=len(chunks), file=file_path)
            return chunks
        except Exception as e:
            log.error("JSON extract error", file=file_path, error=str(e))
            raise ValueError(f"Gagal membaca JSON: {e}")

    # ── Main extraction method ──────────────────────────────────

    def extract(self, file_path: str, extra_metadata: Optional[dict] = None) -> List[ProcessedChunk]:
        """
        Extract teks dari file dan kembalikan sebagai list ProcessedChunk.
        Sync method — panggil via run_in_executor untuk async.
        """
        p = Path(file_path)
        if not p.exists():
            raise FileNotFoundError(f"File tidak ditemukan: {file_path}")

        ext = p.suffix.lower()
        file_meta = self.extract_metadata(file_path)
        if extra_metadata:
            file_meta.update(extra_metadata)

        # Verify file is not empty
        if file_meta["file_size_bytes"] == 0:
            raise ValueError(f"File kosong (0 bytes): {file_path}")

        log.info("Processing document", file=p.name, ext=ext, size_kb=file_meta["file_size_kb"])

        # Route ke extractor yang sesuai
        if ext == ".pdf":
            segments = self._extract_pdf(file_path)
        elif ext == ".docx":
            segments = self._extract_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            segments = self._extract_xlsx(file_path)
        elif ext == ".csv":
            segments = self._extract_csv(file_path)
        elif ext in (".pptx", ".ppt"):
            segments = self._extract_pptx(file_path)
        elif ext == ".json":
            segments = self._extract_json(file_path)
        elif ext in (".txt", ".md"):
            segments = self._extract_text(file_path)
        else:
            raise ValueError(f"Format file tidak didukung: {ext}")

        if not segments:
            log.warning("Tidak ada teks berhasil diekstrak", file=file_path)
            return []

        # Buat chunks dengan metadata
        chunks = []
        for i, segment in enumerate(segments):
            if segment.strip():
                meta = {**file_meta, "segment_index": i, "total_segments": len(segments)}
                chunks.append(ProcessedChunk(content=segment, metadata=meta))

        log.info("Document processing complete", file=p.name, chunks=len(chunks))
        return chunks

    async def extract_async(
        self,
        file_path: str,
        extra_metadata: Optional[dict] = None,
        timeout: int = PROCESS_TIMEOUT,
    ) -> List[ProcessedChunk]:
        """
        Async wrapper dengan timeout untuk extract().
        Timeout default 45 detik.
        """
        loop = asyncio.get_event_loop()
        try:
            chunks = await asyncio.wait_for(
                loop.run_in_executor(None, self.extract, file_path, extra_metadata),
                timeout=timeout,
            )
            return chunks
        except asyncio.TimeoutError:
            log.error(
                "Document processing timeout",
                file=file_path,
                timeout=timeout,
            )
            raise TimeoutError(f"Pemrosesan dokumen melebihi {timeout} detik: {Path(file_path).name}")
        except Exception as e:
            log.error("Document processing failed", file=file_path, error=str(e))
            raise


# Singleton instance
document_processor = DocumentProcessor()
