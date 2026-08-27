import asyncio
import os
import json
import re
import re as _re
import socket
import structlog
import hashlib
from datetime import datetime
from core.model_manager import model_manager

log = structlog.get_logger()

# ── Sandbox integration (graceful — no error if unavailable) ─────────────────
try:
    from core.sandbox_executor import sandbox_executor, SandboxProfile
    _SANDBOX_AVAILABLE = True
except ImportError:
    _SANDBOX_AVAILABLE = False
    sandbox_executor = None

# ── Concurrency & File Locking ──────────────────────────────────────────────
try:
    from core.concurrency import FileLock, resource_semaphore
    _CONCURRENCY_AVAILABLE = True
except ImportError:
    _CONCURRENCY_AVAILABLE = False
    FileLock = None
    resource_semaphore = None


# ─── Foreground Server Detection ────────────────────────────────────────────
FOREGROUND_SERVER_PATTERNS = [
    r'\bnode\s+\S+\.(?:js|mjs|cjs)\b',
    r'\bnpm\s+start\b',
    r'\bnpm\s+run\s+(?:dev|serve|start)\b',
    r'\byarn\s+(?:start|dev)\b',
    r'\bpnpm\s+(?:start|dev)\b',
    r'\bstreamlit\s+run\b',
    r'\buvicorn\b',
    r'\bgunicorn\b',
    r'\bflask\s+run\b',
    r'\bserve\s+-s\b',
    r'\bnpx\s+serve\b',
    r'\bhttp-server\b',
    r'\blive-server\b',
    r'\bphp\s+-S\b',
    r'\bruby\s+\S+\.rb\b',
    r'\bvite\s+(?:preview|dev)?\s*$',
    r'\bnext\s+(?:dev|start)\b',
]

_BACKGROUND_MARKERS = ['&', 'nohup ', '> /dev/null', 'disown', 'setsid']

def _is_foreground_server(command: str) -> bool:
    cmd_stripped = command.strip()
    if any(marker in cmd_stripped for marker in _BACKGROUND_MARKERS):
        return False
    for pattern in FOREGROUND_SERVER_PATTERNS:
        if re.search(pattern, cmd_stripped, re.IGNORECASE):
            return True
    return False

def _wrap_as_background(command: str, project_base_path: str) -> tuple[str, str]:
    import hashlib
    cmd_hash = hashlib.md5(command.encode()).hexdigest()[:8]
    log_file = os.path.join(project_base_path, f".server_{cmd_hash}.log")
    wrapped = f"nohup {command} > {log_file} 2>&1 &"
    return wrapped, log_file

def _classify_command_timeout(command: str) -> float:
    cmd_lower = command.lower().strip()
    if any(p in cmd_lower for p in [
        'npm install', 'npm i ', 'npm ci',
        'pip install', 'pip3 install',
        'apt-get install', 'apt install',
        'yarn install', 'yarn add',
        'pnpm install', 'pnpm add',
        'npx prisma',  # FIX: prisma commands butuh waktu lama
    ]):
        return 300.0
    if any(p in cmd_lower for p in [
        'npm run build', 'yarn build', 'pnpm build',
        'webpack', 'tsc ', 'vite build',
        'next build', 'cargo build', 'make ',
    ]):
        return 180.0
    if cmd_lower.rstrip().endswith('&'):
        return 15.0
    return 60.0


# ─── Port Safety System ─────────────────────────────────────────────────────
RESERVED_PORTS = {7860, 6379, 5432, 3306, 11434}

def _get_safe_port(preferred: int = 0, range_start: int = 8100, range_end: int = 9000) -> int:
    if preferred and preferred not in RESERVED_PORTS:
        try:
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
                s.bind(("127.0.0.1", preferred))
                return preferred
        except OSError:
            pass

    for port in range(range_start, range_end):
        if port in RESERVED_PORTS:
            continue
        try:
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
                s.bind(("127.0.0.1", port))
                return port
        except OSError:
            continue
    return 0

def _rewrite_port_in_command(command: str) -> tuple[str, str]:
    port_patterns = [
        _re.compile(r'(--port[= ]+)(\d+)'),
        _re.compile(r'(-p\s+)(\d+)'),
        _re.compile(r'(PORT=)(\d+)'),
        _re.compile(r'((?:0\.0\.0\.0|127\.0\.0\.1|localhost):)(\d+)'),
        _re.compile(r'(--bind\s+\S+:)(\d+)'),
        _re.compile(r'(run\s+-p\s+)(\d+)'),
        _re.compile(r'(--port\s+)(\d+)'),
    ]
    warning = ""
    modified = command
    for pattern in port_patterns:
        match = pattern.search(modified)
        if match:
            port_str = match.group(2)
            try:
                port = int(port_str)
            except ValueError:
                continue
            if port in RESERVED_PORTS:
                safe_port = _get_safe_port(range_start=port + 1)
                if not safe_port:
                    safe_port = _get_safe_port()
                if safe_port:
                    modified = modified[:match.start(2)] + str(safe_port) + modified[match.end(2):]
                    warning = (
                        f"PORT COLLISION PREVENTED: Port {port} is reserved. "
                        f"Automatically reassigned to port {safe_port}.\n"
                    )
                    break
    return modified, warning


async def find_safe_port(preferred: int = 0) -> str:
    port = _get_safe_port(preferred)
    if port:
        return f"Safe port found: {port}. Free and does NOT conflict with AI Orchestrator (port 7860)."
    return "Error: Could not find a free port in range 8100-9000."


async def execute_bash(command: str, session_id: str = None) -> str:
    """Run a bash command and return output."""
    try:
        # ── Resolve project base path ─────────────────────────────────────
        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                import json as _json
                                meta = _json.loads(meta)
                            except Exception:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if not project_base_path:
            safe_folder = (session_id[:8] if session_id else "agent")
            project_base_path = os.path.expanduser(f"~/projects/{safe_folder}")

        os.makedirs(project_base_path, exist_ok=True)

        # ── FIX: Jika command mengandung 'cd X && ...', jalankan dari home
        # bukan dari project_base_path agar path relatif dalam cd bisa resolve
        # Deteksi apakah command sudah punya 'cd' di awal
        cmd_stripped = command.strip()
        has_explicit_cd = bool(re.match(r'^cd\s+', cmd_stripped))

        # Tentukan cwd yang tepat:
        # - Jika command dimulai dengan 'cd /absolute/path', pakai home sebagai base
        # - Jika command dimulai dengan 'cd relative/path', pakai home sebagai base
        # - Jika tidak ada cd, pakai project_base_path
        if has_explicit_cd:
            # FIX: Jalankan dari home agar cd bisa navigate ke mana saja
            cwd = os.path.expanduser("~")
        else:
            cwd = project_base_path

        # ── Security: Normalize dan cek perintah berbahaya ───────────────
        normalized_cmd = re.sub(r'\s+', ' ', command.lower().strip())

        RESTRICTED_PATTERNS = [
            (r'\brm\s+-rf\s+[/~]', "Destructive rm -rf on root/home"),
            (r'\brm\s+-r\s+[/~]', "Destructive rm -r on root/home"),
            (r'\brmdir\s+/', "rmdir on root"),
            (r'\bmkfs\b', "Disk format command"),
            (r'\bdd\s+if=', "Raw disk write via dd"),
            (r'\b(shutdown|reboot|halt|poweroff)\b', "System power control"),
            (r'\bsystemctl\s+(stop|disable)\s+.*orchestrator', "Stopping AI Orchestrator service"),
            (r':\(\)\s*\{', "Fork bomb pattern"),
            (r':\s*\(\s*\)\s*\{', "Fork bomb pattern"),
            (r'cat\s+.*/\.ssh/', "Reading SSH keys"),
            (r'cat\s+.*/etc/shadow', "Reading shadow passwords"),
            (r'>\s*.*\.env\b', "Overwriting .env"),
            (r'nc\s+.*-e\s+/bin', "Netcat reverse shell"),
            (r'curl\s+.*\|\s*(bash|sh)', "Remote code execution via curl pipe"),
            (r'wget\s+.*\|\s*(bash|sh)', "Remote code execution via wget pipe"),
            (r'\bchmod\s+[0-7]*[67][0-7]\s+/(etc|bin|usr)', "Privilege escalation via chmod"),
            (r'lsof\s+.*:7860', "Querying Orchestrator port 7860"),
            (r'fuser\s+.*7860', "Killing Orchestrator port 7860"),
            (r'\b(pkill|killall)\s+.*(uvicorn|python|orchestrator|main\.py)', "Killing Orchestrator process by name"),
        ]

        for pattern, reason in RESTRICTED_PATTERNS:
            if re.search(pattern, normalized_cmd):
                log.warning("Blocked dangerous command", reason=reason, cmd=command[:80])
                return f"Security Exception: Command blocked — {reason}."

        # Protect Orchestrator PIDs dynamically to prevent suicide
        my_pid = str(os.getpid())
        parent_pid = str(os.getppid())
        if "kill" in normalized_cmd and (my_pid in normalized_cmd or parent_pid in normalized_cmd):
            log.warning("Blocked attempt to kill Orchestrator PID", cmd=command[:80])
            return "Security Exception: Attempting to kill the Orchestrator process is strictly blocked."

        PROTECTED_PATHS = ["/etc/passwd", "/etc/shadow", "/.ssh/id_", "ai-orchestrator.db"]
        cmd_lower = command.lower()
        for protected in PROTECTED_PATHS:
            if protected in cmd_lower and any(op in cmd_lower for op in [" > ", ">>", "tee ", "write"]):
                return f"Security Exception: Writing to protected file '{protected}' is blocked."

        # ── Port collision prevention ──────────────────────────────────────
        port_warning = ""
        command, port_warning = _rewrite_port_in_command(command)

        # ── Auto-background foreground servers ────────────────────────────
        bg_log_file = None
        was_auto_backgrounded = False
        if _is_foreground_server(command):
            log.info("Auto-backgrounding foreground server command", cmd=command[:80])
            command, bg_log_file = _wrap_as_background(command, project_base_path)
            was_auto_backgrounded = True

        # ── Adaptive timeout ──────────────────────────────────────────────
        timeout = _classify_command_timeout(command)

        # FIX: Log command dengan %s bukan f-string untuk hindari format specifier error
        log.info("execute_bash running", cmd=command[:120], cwd=cwd, timeout=timeout)

        # ── Sandbox audit: log profile detection (non-blocking) ──────────
        if _SANDBOX_AVAILABLE and sandbox_executor:
            detected_profile = sandbox_executor.detect_profile(command)
            log.debug("Sandbox profile detected",
                      profile=detected_profile.value, cmd=command[:60])

        proc = await asyncio.create_subprocess_shell(
            command,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=cwd,   # FIX: gunakan cwd yang sudah dihitung di atas
        )

        async def read_stream(stream, max_lines=10000):
            lines = []
            truncated = False
            if stream is None:
                return "", False
            async for line in stream:
                if len(lines) < max_lines:
                    lines.append(line.decode(errors="replace"))
                else:
                    if not truncated:
                        truncated = True
                        try:
                            proc.terminate()
                        except Exception:
                            pass
            return "".join(lines), truncated

        async def run_with_timeout():
            stdout_task = asyncio.create_task(read_stream(proc.stdout))
            stderr_task = asyncio.create_task(read_stream(proc.stderr))
            await proc.wait()
            out, out_trunc = await stdout_task
            err, err_trunc = await stderr_task
            return out, out_trunc, err, err_trunc

        out, out_trunc, err, err_trunc = await asyncio.wait_for(
            run_with_timeout(), timeout=timeout
        )

        output = f"[Working Directory: {cwd}]\n"
        if port_warning:
            output += port_warning

        if was_auto_backgrounded and bg_log_file:
            await asyncio.sleep(3)
            try:
                if os.path.exists(bg_log_file):
                    with open(bg_log_file, 'r', errors='replace') as f:
                        startup_log = f.read(4000)
                    output += "Server started in background (auto-detected foreground command).\n"
                    output += f"Log file: {bg_log_file}\n"
                    if startup_log.strip():
                        output += f"--- Startup Log ---\n{startup_log}\n---\n"
                    else:
                        output += "(Server starting, no log output yet)\n"
                else:
                    output += f"Server started in background. Log file: {bg_log_file}\n"
            except Exception as log_err:
                output += f"Server started in background. (Could not read log: {log_err})\n"
        else:
            if out:
                output += out
                if out_trunc:
                    output += "\n...[STDOUT TRUNCATED]..."
            if err:
                output += "\n[stderr]\n" + err
                if err_trunc:
                    output += "\n...[STDERR TRUNCATED]..."

        result = output.strip() or "Command executed successfully with no output."

        # FIX: Tambahkan hint berguna jika npm/prisma gagal karena package.json tidak ada
        if "missing script" in result.lower() or "no such file" in result.lower():
            result += (
                "\n\n[HINT: Pastikan Anda sudah 'cd' ke direktori project yang benar "
                "sebelum menjalankan npm/prisma. Gunakan: cd /path/ke/project && npm install]"
            )
        if "enoent" in result.lower() and "package.json" in result.lower():
            result += (
                "\n\n[HINT: package.json tidak ditemukan. "
                "Buat dulu dengan write_file atau pastikan path project sudah benar.]"
            )

        return result

    except asyncio.TimeoutError:
        try:
            proc.kill()
        except Exception:
            pass
        return (
            f"Error: Command timed out after {int(timeout)} seconds. "
            "If this is a server command, run it in background: "
            "nohup <command> > app.log 2>&1 &"
        )
    except Exception as e:
        # FIX: Gunakan str(e) bukan f-string dengan format specifier
        return "Error executing command: " + str(e)


async def read_file(path: str, session_id: str = None) -> str:
    """Read a file."""
    try:
        import aiofiles

        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                import json as _json
                                meta = _json.loads(meta)
                            except Exception:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if project_base_path and not os.path.isabs(path):
            abs_path = os.path.join(project_base_path, path)
        elif not os.path.isabs(path) and session_id:
            safe_folder = session_id[:8]
            abs_path = os.path.join(os.path.expanduser(f"~/projects/{safe_folder}"), path)
        else:
            abs_path = os.path.abspath(path)

        abs_real = os.path.realpath(abs_path)

        _BLOCKED_READ_PATTERNS = [
            "/etc/shadow", "/etc/passwd", "/etc/sudoers",
            "/.ssh/", "/.gnupg/", "/proc/", "/sys/",
        ]
        _BLOCKED_READ_EXTENSIONS = [".pem", ".key", ".p12", ".pfx"]
        _is_blocked = any(p in abs_real for p in _BLOCKED_READ_PATTERNS)
        _is_blocked = _is_blocked or any(abs_real.endswith(ext) for ext in _BLOCKED_READ_EXTENSIONS)

        # Blokir baca file .env milik orchestrator sendiri
        _env_path = os.path.realpath(
            os.path.join(os.path.dirname(__file__), "..", "..", "..", ".env")
        )
        if abs_real == _env_path:
            _is_blocked = True

        if _is_blocked:
            log.warning("Blocked read_file on sensitive path", path=abs_real[:100])
            return f"Security Exception: Reading '{os.path.basename(abs_real)}' is blocked."

        if not os.path.exists(abs_path):
            return f"Error: File not found: {abs_path}"

        async def _read():
            async with aiofiles.open(abs_path, "r", encoding="utf-8") as f:
                return await f.read()

        return await asyncio.wait_for(_read(), timeout=30.0)

    except asyncio.TimeoutError:
        return f"Error: Reading file {path} timed out after 30 seconds."
async def _update_artifact_registry(session_id: str, abs_path: str, content: str):
    """Helper to track file changes and their hashes."""
    if not session_id:
        return
    try:
        from db.database import AsyncSessionLocal
        from db.models import ArtifactRegistry
        from sqlmodel import select
        
        file_hash = hashlib.sha256(content.encode()).hexdigest()
        
        async with AsyncSessionLocal() as db:
            result = await db.execute(
                select(ArtifactRegistry).where(
                    ArtifactRegistry.session_id == session_id,
                    ArtifactRegistry.file_path == abs_path
                )
            )
            existing = result.scalars().first()
            if existing:
                existing.file_hash = file_hash
                existing.last_modified = datetime.now().replace(tzinfo=None)
                db.add(existing)
            else:
                new_entry = ArtifactRegistry(
                    session_id=session_id,
                    file_path=abs_path,
                    file_hash=file_hash
                )
                db.add(new_entry)
            await db.commit()
    except Exception as e:
        log.debug("Artifact Registry update skipped", error=str(e))


async def write_file(path: str, content: str, session_id: str = None, confirm: bool = False) -> str:
    """Write content to a file with safety checks for overwrites."""
    try:
        import aiofiles

        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                import json as _json
                                meta = _json.loads(meta)
                            except Exception:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if not project_base_path and not os.path.isabs(path):
            safe_folder = session_id[:8] if session_id else "agent"
            project_base_path = os.path.expanduser(f"~/projects/{safe_folder}")

        if project_base_path and not os.path.isabs(path):
            path = os.path.join(project_base_path, path)

        abs_path = os.path.abspath(path)

        # ── Safety Checks ──────────────────────────────────────────────────
        projects_base = os.path.expanduser("~/projects")
        is_outside = not abs_path.startswith(projects_base)
        
        if os.path.exists(abs_path):
            size = os.path.getsize(abs_path)
            # Require confirmation for large overwrites or writing outside projects
            if (size > 100 * 1024 or is_outside) and not confirm:
                return (
                    f"⚠️ File '{abs_path}' sudah ada "
                    f"({size/1024:.1f}KB" + (" dan di luar workspace" if is_outside else "") + ").\n"
                    f"   Gunakan confirm=True untuk menimpa file ini secara paksa."
                )

        os.makedirs(os.path.dirname(abs_path), exist_ok=True)

        async def _write():
            if _CONCURRENCY_AVAILABLE and FileLock and resource_semaphore:
                async with resource_semaphore.acquire("filesystem", max_concurrent=5):
                    async with FileLock(abs_path, timeout=15.0):
                        async with aiofiles.open(abs_path, "w", encoding="utf-8") as f:
                            await f.write(content)
            else:
                async with aiofiles.open(abs_path, "w", encoding="utf-8") as f:
                    await f.write(content)

        await asyncio.wait_for(_write(), timeout=30.0)

        # Update Artifact Registry
        await _update_artifact_registry(session_id, abs_path, content)

        display_path = os.path.relpath(abs_path, project_base_path) if project_base_path else path
        result = "Successfully wrote to " + display_path
        if project_base_path:
            result = "Project: " + project_base_path + "\n" + result
            try:
                from core.project_indexer import project_indexer
                asyncio.create_task(project_indexer.scan_project(session_id or "", project_base_path))
            except Exception:
                pass

        return result

    except asyncio.TimeoutError:
        return "Error: Writing file " + path + " timed out after 30 seconds."
    except Exception as e:
        return "Error writing file " + path + ": " + str(e)


async def write_multiple_files(files_data: list, session_id: str = None) -> str:
    """Write multiple files at once."""
    try:
        import aiofiles

        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                import json as _json
                                meta = _json.loads(meta)
                            except Exception:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if not project_base_path:
            first_relative = next(
                (f.get("path", "") for f in files_data
                 if f.get("path") and not os.path.isabs(f.get("path", ""))),
                None
            )
            if first_relative:
                safe_folder = session_id[:8] if session_id else "agent"
                project_base_path = os.path.expanduser(f"~/projects/{safe_folder}")

        results = []
        for file_obj in files_data:
            path = file_obj.get("path")
            content = file_obj.get("content", "")
            if not path:
                results.append("Skipped invalid entry: " + str(file_obj))
                continue

            if project_base_path and not os.path.isabs(path):
                path = os.path.join(project_base_path, path)

            try:
                abs_path = os.path.abspath(path)
                
                # ── Safety Checks ──────────────────────────────────────────
                projects_base = os.path.expanduser("~/projects")
                is_outside = not abs_path.startswith(projects_base)
                
                if os.path.exists(abs_path):
                    size = os.path.getsize(abs_path)
                    if (size > 100 * 1024 or is_outside) and not confirm:
                        results.append(
                            f"Skipped {path}: File exists ({size/1024:.1f}KB" +
                            (" and outside workspace" if is_outside else "") + "). Needs confirm=True."
                        )
                        continue

                os.makedirs(os.path.dirname(abs_path), exist_ok=True)
                async with aiofiles.open(abs_path, "w", encoding="utf-8") as f:
                    await f.write(content)
                display_path = os.path.relpath(abs_path, project_base_path) if project_base_path else path
                results.append("OK: " + display_path)
                # Update Artifact Registry
                await _update_artifact_registry(session_id, abs_path, content)
            except Exception as e:
                results.append("Error writing " + path + ": " + str(e))

        if project_base_path:
            try:
                from core.project_indexer import project_indexer
                asyncio.create_task(project_indexer.scan_project(session_id or "", project_base_path))
            except Exception:
                pass
            return "Project: " + project_base_path + "\n" + "\n".join(results)
        else:
            return "\n".join(results)

    except Exception as e:
        return "Fatal error in write_multiple_files: " + str(e)


async def write_file_chunk(path: str, content: str, chunk_index: int, total_chunks: int, session_id: str = None) -> str:
    """
    Write file in chunks for large files or truncation recovery.
    Appends content if chunk_index > 0.
    """
    try:
        import aiofiles
        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                import json as _json
                                meta = _json.loads(meta)
                            except:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if not project_base_path and not os.path.isabs(path):
            safe_folder = session_id[:8] if session_id else "agent"
            project_base_path = os.path.expanduser(f"~/projects/{safe_folder}")

        if project_base_path and not os.path.isabs(path):
            path = os.path.join(project_base_path, path)

        abs_path = os.path.abspath(path)
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)

        mode = "w" if chunk_index == 0 else "a"
        async with aiofiles.open(abs_path, mode, encoding="utf-8") as f:
            await f.write(content)

        if chunk_index + 1 == total_chunks:
            # Final chunk, update registry
            with open(abs_path, 'r', encoding='utf-8') as f:
                full_content = f.read()
            await _update_artifact_registry(session_id, abs_path, full_content)
            
            try:
                from core.project_indexer import project_indexer
                asyncio.create_task(project_indexer.scan_project(session_id or "", project_base_path or os.path.dirname(abs_path)))
            except Exception:
                pass

        return f"Successfully wrote chunk {chunk_index + 1}/{total_chunks} to {os.path.basename(abs_path)}"
    except Exception as e:
        return f"Error writing chunk: {str(e)}"


async def ask_model(model_id: str, prompt: str) -> str:
    """Ask another AI model a question."""
    try:
        if model_id not in model_manager.available_models:
            return "Error: Model '" + model_id + "' is not available."
        messages = [{"role": "user", "content": prompt}]
        response = ""
        async for chunk in model_manager.chat_stream(
            model_id, messages, temperature=0.7, max_tokens=4096
        ):
            if chunk and not chunk.startswith("\n[Error"):
                response += chunk
        return response.strip() or "No response from model."
    except Exception as e:
        return "Error asking model " + model_id + ": " + str(e)

async def list_dir(path: str, session_id: str = None) -> str:
    """List directory contents."""
    try:
        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                import json as _json
                                meta = _json.loads(meta)
                            except:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if project_base_path and not os.path.isabs(path):
            abs_path = os.path.join(project_base_path, path)
        elif not os.path.isabs(path) and session_id:
            safe_folder = session_id[:8]
            abs_path = os.path.join(os.path.expanduser(f"~/projects/{safe_folder}"), path)
        else:
            abs_path = os.path.abspath(path)

        if not os.path.exists(abs_path) or not os.path.isdir(abs_path):
            return f"Error: Directory not found or is not a directory: {abs_path}"

        items = os.listdir(abs_path)
        result = []
        for item in sorted(items):
            full_item = os.path.join(abs_path, item)
            if os.path.isdir(full_item):
                result.append(f"[DIR]  {item}/")
            else:
                size = os.path.getsize(full_item)
                result.append(f"[FILE] {item} ({size} bytes)")

        return f"Contents of {abs_path}:\n" + "\n".join(result)
    except Exception as e:
        return "Error listing directory: " + str(e)

async def search_files(query: str, path: str = ".", session_id: str = None) -> str:
    """Search for files matching query or content."""
    try:
        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                import json as _json
                                meta = _json.loads(meta)
                            except:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if project_base_path and not os.path.isabs(path):
            abs_path = os.path.join(project_base_path, path)
        elif not os.path.isabs(path) and session_id:
            safe_folder = session_id[:8]
            abs_path = os.path.join(os.path.expanduser(f"~/projects/{safe_folder}"), path)
        else:
            abs_path = os.path.abspath(path)

        if not os.path.exists(abs_path):
            return f"Error: Path not found: {abs_path}"

        import shlex
        safe_query = shlex.quote(query)
        safe_path = shlex.quote(abs_path)
        
        # 1. Cari konten
        cmd_content = f"grep -irl {safe_query} {safe_path} --exclude-dir=node_modules --exclude-dir=.git --exclude-dir=venv | head -n 20"
        proc = await asyncio.create_subprocess_shell(cmd_content, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
        stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30.0)
        output = stdout.decode().strip()

        # 2. Jika tidak ada konten, cari nama file
        if not output:
            cmd_name = f"find {safe_path} -type f -not -path '*/node_modules/*' -not -path '*/.git/*' -iname '*{query}*' | head -n 20"
            proc2 = await asyncio.create_subprocess_shell(cmd_name, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
            stdout2, _ = await asyncio.wait_for(proc2.communicate(), timeout=30.0)
            output = stdout2.decode().strip()

        if not output:
            return f"No results found for '{query}' in {abs_path}"
        return f"Search results for '{query}' in {abs_path}:\n{output}"
    except asyncio.TimeoutError:
        return "Error: Search timed out after 30 seconds."
    except Exception as e:
        return "Error searching files: " + str(e)



async def read_document(path: str, session_id: str = None) -> str:
    """Read a document (PDF, DOCX, XLSX, PPTX, CSV, TXT) seamlessly."""
    try:
        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        import json as _json
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                meta = _json.loads(meta)
                            except Exception:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if project_base_path and not os.path.isabs(path):
            abs_path = os.path.join(project_base_path, path)
        elif not os.path.isabs(path) and session_id:
            safe_folder = session_id[:8]
            abs_path = os.path.join(os.path.expanduser(f"~/projects/{safe_folder}"), path)
        else:
            abs_path = os.path.abspath(path)

        if not os.path.exists(abs_path):
            return f"Error: File not found: {abs_path}"

        # Delegate to RAG document processor
        from rag.document_processor import DocumentProcessor
        processor = DocumentProcessor()
        segments = processor.extract_text(abs_path)
        
        full_text = "\n\n".join(segments)
        
        # Truncate if too long to prevent blowing up the context window
        if len(full_text) > 20000:
            full_text = full_text[:20000] + "\n\n...[DOCUMENT TRUNCATED DUE TO LENGTH]..."
            
        return f"--- Content of {os.path.basename(abs_path)} ---\n{full_text}"

    except Exception as e:
        return f"Error reading document: {str(e)}"


async def replace_in_file(path: str, old_string: str, new_string: str, session_id: str = None) -> str:
    """Surgically replace text in a file."""
    try:
        import aiofiles
        
        project_base_path = None
        if session_id:
            try:
                from db.database import AsyncSessionLocal
                from db.models import ChatSession
                async with AsyncSessionLocal() as db:
                    session = await db.get(ChatSession, session_id)
                    if session and session.project_metadata:
                        import json as _json
                        meta = session.project_metadata
                        if isinstance(meta, str):
                            try:
                                meta = _json.loads(meta)
                            except Exception:
                                meta = {}
                        project_base_path = meta.get("project_path") if isinstance(meta, dict) else None
            except Exception:
                pass

        if project_base_path and not os.path.isabs(path):
            abs_path = os.path.join(project_base_path, path)
        elif not os.path.isabs(path) and session_id:
            safe_folder = session_id[:8]
            abs_path = os.path.join(os.path.expanduser(f"~/projects/{safe_folder}"), path)
        else:
            abs_path = os.path.abspath(path)

        if not os.path.exists(abs_path):
            return f"Error: File not found: {abs_path}"

        async with aiofiles.open(abs_path, "r", encoding="utf-8") as f:
            content = await f.read()

        if old_string not in content:
            # Maybe it has different line endings, try to normalize
            normalized_content = content.replace("\r\n", "\n")
            normalized_old = old_string.replace("\r\n", "\n")
            if normalized_old not in normalized_content:
                return f"Error: The exact old_string was not found in {path}. Make sure to include the exact literal text including spaces and newlines."
            else:
                content = normalized_content
                old_string = normalized_old
                new_string = new_string.replace("\r\n", "\n")

        count = content.count(old_string)
        if count > 1:
            return f"Error: old_string is ambiguous. Found {count} occurrences in {path}. Please provide a larger snippet to uniquely identify the section to replace."

        new_content = content.replace(old_string, new_string)

        async with aiofiles.open(abs_path, "w", encoding="utf-8") as f:
            await f.write(new_content)
            
        await _update_artifact_registry(session_id, abs_path, new_content)

        return f"Successfully replaced {len(old_string)} chars with {len(new_string)} chars in {path}."

    except Exception as e:
        return f"Error replacing text in file: {str(e)}"


# ─── Proactive Scheduler Tool ─────────────────────────────────────────────────

async def schedule_task(
    title: str,
    description: str,
    due_in_minutes: int = 60,
    recurrence: str = None,
    session_id: str = None,
    user_id: str = None,
) -> str:
    """
    Jadwalkan tugas atau pengingat proaktif yang akan dieksekusi di masa depan.
    
    Args:
        title: Judul singkat tugas (maks 100 karakter)
        description: Instruksi lengkap yang akan dijalankan saat due_at tercapai
        due_in_minutes: Berapa menit dari sekarang task akan dieksekusi (default: 60)
        recurrence: Pola perulangan — None, "daily", "weekly", atau cron expression
        session_id: ID sesi saat ini (diisi otomatis)
        user_id: ID user (diisi otomatis)
    
    Returns:
        Konfirmasi penjadwalan task beserta ID-nya
    """
    try:
        from datetime import timezone, timedelta
        from db.database import SessionLocal
        from db.models import ScheduledTask

        due_at = datetime.now(timezone.utc).replace(tzinfo=None) + timedelta(minutes=int(due_in_minutes))

        with SessionLocal() as db:
            task = ScheduledTask(
                session_id=session_id or "unknown",
                user_id=user_id or "unknown",
                title=str(title)[:200],
                description=str(description)[:2000],
                due_at=due_at,
                recurrence=str(recurrence) if recurrence else None,
                status="pending",
            )
            db.add(task)
            db.commit()
            db.refresh(task)
            task_id = task.id

        due_str = due_at.strftime("%Y-%m-%d %H:%M UTC")
        rec_str = f" (berulang: {recurrence})" if recurrence else ""
        log.info("Scheduled proactive task", task_id=task_id, due_at=due_str, title=title)
        return (
            f"✅ Task terjadwal berhasil dibuat!\n"
            f"ID: {task_id}\n"
            f"Judul: {title}\n"
            f"Akan dijalankan: {due_str}{rec_str}\n"
            f"Instruksi: {str(description)[:200]}..."
        )
    except Exception as e:
        log.error("Failed to schedule task", error=str(e))
        return f"Error menjadwalkan task: {str(e)}"


# ═══════════════════════════════════════════════════════════════════════════════
# OFFICE & COMPUTER USE TOOLS (v5.0)
# Menggunakan library yang sudah ada: openpyxl, python-docx, python-pptx
# ═══════════════════════════════════════════════════════════════════════════════

# ── Excel Tools ───────────────────────────────────────────────────────────────

async def excel_read(path: str, sheet: str = None, max_rows: int = 100) -> str:
    """
    Baca file Excel (.xlsx) dan kembalikan sebagai teks tabel.

    Args:
        path: Path ke file Excel
        sheet: Nama sheet (default: sheet pertama)
        max_rows: Maksimal baris yang dibaca (default: 100)

    Returns:
        Isi Excel sebagai teks tabel
    """
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        ws = wb[sheet] if sheet else wb.active
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                rows.append(f"... (dan {ws.max_row - max_rows} baris lagi)")
                break
            row_str = " | ".join(str(c) if c is not None else "" for c in row)
            rows.append(row_str)
        wb.close()
        sheet_name = ws.title
        return f"📊 Excel: {path} [Sheet: {sheet_name}] ({ws.max_row} baris, {ws.max_column} kolom)\n\n" + "\n".join(rows)
    except ImportError:
        return "Error: openpyxl tidak terinstall. Jalankan: pip install openpyxl"
    except Exception as e:
        return f"Error membaca Excel: {e}"


async def excel_write(
    path: str,
    data: list,
    sheet: str = "Sheet1",
    headers: list = None,
    overwrite: bool = False
) -> str:
    """
    Tulis data ke file Excel (.xlsx).

    Args:
        path: Path output file Excel
        data: List of lists (baris data), contoh: [["Alice", 25], ["Bob", 30]]
        sheet: Nama sheet
        headers: Nama kolom, contoh: ["Nama", "Umur"]
        overwrite: True untuk menimpa file yang ada

    Returns:
        Konfirmasi penulisan
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        if os.path.exists(path) and not overwrite:
            return f"File {path} sudah ada. Set overwrite=true untuk menimpa."

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet

        # Tulis headers dengan styling
        if headers:
            ws.append(headers)
            header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            for cell in ws[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")

        # Tulis data
        row_count = 0
        if isinstance(data, list):
            for row in data:
                if isinstance(row, list):
                    ws.append(row)
                elif isinstance(row, dict):
                    ws.append(list(row.values()))
                row_count += 1

        # Auto-fit kolom
        for col in ws.columns:
            max_len = max((len(str(c.value)) for c in col if c.value), default=10)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

        os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
        wb.save(path)
        wb.close()
        return f"✅ Excel berhasil dibuat: {path}\n📊 Sheet: {sheet} | {row_count} baris data" + (f" | {len(headers)} kolom" if headers else "")
    except Exception as e:
        return f"Error menulis Excel: {e}"


async def excel_formula(path: str, cell: str, formula: str, sheet: str = None) -> str:
    """
    Masukkan formula ke sel Excel tertentu.

    Args:
        path: Path ke file Excel
        cell: Sel target, contoh: "D2" atau "SUM_ROW" untuk auto-sum
        formula: Formula Excel, contoh: "=SUM(B2:B10)" atau "=AVERAGE(C2:C10)"
        sheet: Nama sheet (default: aktif)

    Returns:
        Konfirmasi
    """
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet else wb.active
        ws[cell] = formula
        wb.save(path)
        wb.close()
        return f"✅ Formula '{formula}' dimasukkan ke sel {cell} pada {path}"
    except Exception as e:
        return f"Error: {e}"


# ── Word Tools ────────────────────────────────────────────────────────────────

async def word_read(path: str) -> str:
    """
    Baca file Word (.docx) dan kembalikan isinya sebagai teks.

    Args:
        path: Path ke file .docx

    Returns:
        Isi dokumen Word sebagai teks
    """
    try:
        from docx import Document
        doc = Document(path)
        sections = []

        # Baca paragraf
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "")
                prefix = "#" * int(level) if level.isdigit() else "#"
                sections.append(f"{prefix} {para.text}")
            elif para.text.strip():
                sections.append(para.text)

        # Baca tabel jika ada
        for i, table in enumerate(doc.tables):
            sections.append(f"\n[Tabel {i+1}]")
            for row in table.rows:
                sections.append(" | ".join(cell.text for cell in row.cells))

        return f"📄 Word: {path}\n\n" + "\n".join(sections)
    except ImportError:
        return "Error: python-docx tidak terinstall. Jalankan: pip install python-docx"
    except Exception as e:
        return f"Error membaca Word: {e}"


async def word_write(
    path: str,
    content: str,
    title: str = None,
    author: str = "AI ORCHESTRATOR",
    overwrite: bool = False
) -> str:
    """
    Buat file Word (.docx) dari teks markdown-style.

    Args:
        path: Path output file .docx
        content: Konten dokumen. Gunakan # untuk Heading 1, ## untuk Heading 2, dsb.
        title: Judul dokumen (opsional)
        author: Nama penulis
        overwrite: True untuk menimpa file yang ada

    Returns:
        Konfirmasi pembuatan
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        if os.path.exists(path) and not overwrite:
            return f"File {path} sudah ada. Set overwrite=true untuk menimpa."

        doc = Document()

        # Set properties
        doc.core_properties.author = author
        if title:
            doc.core_properties.title = title
            heading = doc.add_heading(title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Parse konten markdown-style
        lines = content.split("\n")
        i = 0
        while i < len(lines):
            line = lines[i]

            if line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("- ") or line.startswith("* "):
                # Bullet list
                para = doc.add_paragraph(line[2:], style="List Bullet")
            elif line.startswith("1. ") or (len(line) > 2 and line[0].isdigit() and line[1] == "."):
                # Numbered list
                text = line.split(". ", 1)[1] if ". " in line else line
                doc.add_paragraph(text, style="List Number")
            elif line.strip() == "---" or line.strip() == "===":
                doc.add_page_break()
            elif line.strip():
                doc.add_paragraph(line)
            else:
                # Baris kosong → spasi
                pass
            i += 1

        os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
        doc.save(path)
        word_count = len(content.split())
        return f"✅ Word berhasil dibuat: {path}\n📄 {word_count} kata | {len(lines)} baris"
    except Exception as e:
        return f"Error membuat Word: {e}"


# ── PowerPoint Tools ──────────────────────────────────────────────────────────

async def ppt_create(
    path: str,
    title: str,
    subtitle: str = "",
    theme: str = "dark",
    overwrite: bool = False
) -> str:
    """
    Buat file PowerPoint (.pptx) baru dengan slide judul.

    Args:
        path: Path output file .pptx
        title: Judul presentasi
        subtitle: Sub-judul (opsional)
        theme: "dark" atau "light"
        overwrite: True untuk menimpa

    Returns:
        Konfirmasi pembuatan
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        if os.path.exists(path) and not overwrite:
            return f"File {path} sudah ada. Set overwrite=true untuk menimpa."

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide judul
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Background color
        bg_color = RGBColor(0x0F, 0x17, 0x2A) if theme == "dark" else RGBColor(0xFF, 0xFF, 0xFF)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        tf = title_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_run = tf.paragraphs[0].runs[0]
        title_run.font.size = Pt(44)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if theme == "dark" else RGBColor(0x0F, 0x17, 0x2A)

        if subtitle:
            subtitle_shape.text = subtitle
            sub_tf = subtitle_shape.text_frame
            sub_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            sub_run = sub_tf.paragraphs[0].runs[0]
            sub_run.font.size = Pt(24)
            sub_run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8) if theme == "dark" else RGBColor(0x64, 0x74, 0x8B)

        os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
        prs.save(path)
        return f"✅ PowerPoint dibuat: {path}\n📊 Judul: '{title}' | Theme: {theme}"
    except Exception as e:
        return f"Error membuat PowerPoint: {e}"


async def ppt_add_slide(
    path: str,
    title: str,
    content: str,
    layout: str = "content",
    slide_number: int = None
) -> str:
    """
    Tambah slide baru ke file PowerPoint yang ada.

    Args:
        path: Path file .pptx yang sudah ada
        title: Judul slide
        content: Konten slide (teks, bisa berisi bullet poin dengan -)
        layout: "content" (judul+konten), "blank", "title_only"
        slide_number: Posisi slide (default: di akhir)

    Returns:
        Konfirmasi penambahan slide
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        if not os.path.exists(path):
            return f"File {path} tidak ditemukan. Buat dulu dengan ppt_create."

        prs = Presentation(path)

        layout_map = {"content": 1, "blank": 6, "title_only": 5}
        layout_idx = layout_map.get(layout, 1)
        slide_layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
        slide = prs.slides.add_slide(slide_layout)

        # Tambah judul
        if slide.shapes.title:
            slide.shapes.title.text = title
            run = slide.shapes.title.text_frame.paragraphs[0].runs
            if run:
                run[0].font.size = Pt(32)
                run[0].font.bold = True

        # Tambah konten
        if len(slide.placeholders) > 1 and content:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            lines = content.strip().split("\n")
            for j, line in enumerate(lines):
                if j == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()

                line = line.strip()
                if line.startswith("- ") or line.startswith("* "):
                    para.text = line[2:]
                    para.level = 1
                else:
                    para.text = line
                    para.level = 0

                if para.runs:
                    para.runs[0].font.size = Pt(18)

        prs.save(path)
        slide_count = len(prs.slides)
        return f"✅ Slide '{title}' ditambahkan ke {path}\n📊 Total slides: {slide_count}"
    except Exception as e:
        return f"Error menambah slide: {e}"


# ── SQL Tools ─────────────────────────────────────────────────────────────────

async def sql_query(database_path: str, query: str, max_rows: int = 50) -> str:
    """
    Jalankan query SELECT ke database SQLite.

    Args:
        database_path: Path ke file .db SQLite
        query: Query SQL (hanya SELECT yang diizinkan)
        max_rows: Maksimal baris yang dikembalikan

    Returns:
        Hasil query sebagai tabel teks
    """
    try:
        import aiosqlite
        import re as _re

        # Safety: hanya izinkan SELECT
        stripped = query.strip().upper()
        if not stripped.startswith("SELECT") and not stripped.startswith("PRAGMA") and not stripped.startswith("WITH"):
            return "Error: Hanya query SELECT/PRAGMA/WITH yang diizinkan di sql_query. Gunakan sql_execute untuk DDL/DML."

        async with aiosqlite.connect(database_path) as db:
            db.row_factory = aiosqlite.Row
            async with db.execute(query) as cursor:
                rows = await cursor.fetchmany(max_rows)
                if not rows:
                    return f"Query berhasil tapi tidak ada hasil.\nQuery: {query}"

                # Format sebagai tabel
                headers = rows[0].keys()
                header_line = " | ".join(str(h) for h in headers)
                separator = "-+-".join("-" * min(len(str(h)), 20) for h in headers)
                data_lines = []
                for row in rows:
                    data_lines.append(" | ".join(str(v)[:20] if v is not None else "NULL" for v in row))

                total = len(rows)
                result = f"🗄️ Database: {database_path}\nQuery: {query}\n\n"
                result += header_line + "\n" + separator + "\n"
                result += "\n".join(data_lines)
                if total >= max_rows:
                    result += f"\n\n... (ditampilkan {max_rows} baris pertama)"
                return result
    except Exception as e:
        return f"Error SQL query: {e}"


async def sql_execute(database_path: str, statement: str) -> str:
    """
    Jalankan statement DDL/DML ke database SQLite (CREATE, INSERT, UPDATE, DELETE, DROP).

    Args:
        database_path: Path ke file .db SQLite (akan dibuat jika belum ada)
        statement: SQL statement

    Returns:
        Konfirmasi eksekusi dan jumlah baris yang terpengaruh
    """
    try:
        import aiosqlite

        async with aiosqlite.connect(database_path) as db:
            async with db.execute(statement) as cursor:
                await db.commit()
                rows_affected = cursor.rowcount
                lastrowid = cursor.lastrowid

        result = f"✅ SQL dieksekusi pada {database_path}\n"
        result += f"Statement: {statement[:100]}...\n" if len(statement) > 100 else f"Statement: {statement}\n"
        if rows_affected and rows_affected > 0:
            result += f"Baris terpengaruh: {rows_affected}"
        if lastrowid:
            result += f"\nID terakhir yang disisipkan: {lastrowid}"
        return result
    except Exception as e:
        return f"Error SQL execute: {e}"


# ── Screenshot Tool ───────────────────────────────────────────────────────────

async def screenshot_screen(output_path: str = None, region: str = "full") -> str:
    """
    Ambil screenshot layar desktop.

    Args:
        output_path: Path untuk menyimpan screenshot (default: /tmp/screenshot_<ts>.png)
        region: "full" untuk layar penuh, atau "active" untuk jendela aktif

    Returns:
        Path file screenshot yang disimpan
    """
    try:
        import subprocess
        import tempfile

        if not output_path:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"/tmp/screenshot_{ts}.png"

        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else "/tmp", exist_ok=True)

        # Coba berbagai tool screenshot yang tersedia di Linux
        tools_tried = []

        # 1. scrot (ringan, banyak tersedia)
        try:
            result = subprocess.run(
                ["scrot", output_path],
                capture_output=True, timeout=10
            )
            if result.returncode == 0:
                size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
                return f"✅ Screenshot disimpan: {output_path} ({size // 1024}KB)"
            tools_tried.append("scrot")
        except FileNotFoundError:
            tools_tried.append("scrot (tidak ada)")
        except Exception:
            pass

        # 2. gnome-screenshot
        try:
            result = subprocess.run(
                ["gnome-screenshot", "-f", output_path],
                capture_output=True, timeout=10
            )
            if result.returncode == 0:
                return f"✅ Screenshot disimpan: {output_path}"
            tools_tried.append("gnome-screenshot")
        except FileNotFoundError:
            tools_tried.append("gnome-screenshot (tidak ada)")
        except Exception:
            pass

        # 3. import (ImageMagick)
        try:
            result = subprocess.run(
                ["import", "-window", "root", output_path],
                capture_output=True, timeout=10
            )
            if result.returncode == 0:
                return f"✅ Screenshot disimpan: {output_path}"
            tools_tried.append("import/ImageMagick")
        except FileNotFoundError:
            tools_tried.append("import/ImageMagick (tidak ada)")
        except Exception:
            pass

        # 4. xwd + convert
        try:
            xwd_path = output_path.replace(".png", ".xwd")
            r1 = subprocess.run(["xwd", "-root", "-out", xwd_path], capture_output=True, timeout=10)
            if r1.returncode == 0:
                r2 = subprocess.run(["convert", xwd_path, output_path], capture_output=True, timeout=10)
                os.remove(xwd_path)
                if r2.returncode == 0:
                    return f"✅ Screenshot disimpan: {output_path}"
        except Exception:
            pass

        return f"⚠️ Screenshot gagal. Tools dicoba: {', '.join(tools_tried)}\nInstall dengan: sudo apt install scrot"

    except Exception as e:
        return f"Error screenshot: {e}"


# ── App Plan Generator ─────────────────────────────────────────────────────────

async def app_plan_generate(
    description: str,
    app_type: str = "web",
    output_path: str = None
) -> str:
    """
    Generate rencana aplikasi terstruktur dari deskripsi natural language.

    Args:
        description: Deskripsi aplikasi yang ingin dibuat
        app_type: "web", "mobile", "desktop", "api", "cli"
        output_path: Path untuk menyimpan plan JSON (opsional)

    Returns:
        Rencana aplikasi dalam format terstruktur
    """
    try:
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")

        # Template rencana berdasarkan app_type
        templates = {
            "web": {
                "tech_stack": {
                    "frontend": "React + Vite + TailwindCSS",
                    "backend": "FastAPI (Python)",
                    "database": "SQLite / PostgreSQL",
                    "deployment": "Docker + nginx"
                },
                "folder_structure": [
                    "frontend/src/components/",
                    "frontend/src/pages/",
                    "frontend/src/hooks/",
                    "frontend/src/utils/",
                    "backend/api/",
                    "backend/core/",
                    "backend/models/",
                    "docker-compose.yml",
                    ".env.example",
                    "README.md"
                ]
            },
            "mobile": {
                "tech_stack": {
                    "framework": "Flutter (Dart)",
                    "state_management": "Riverpod / Bloc",
                    "backend": "FastAPI (Python)",
                    "database": "SQLite (lokal) / Supabase (cloud)"
                },
                "folder_structure": [
                    "lib/screens/",
                    "lib/widgets/",
                    "lib/providers/",
                    "lib/models/",
                    "lib/services/",
                    "assets/",
                    "pubspec.yaml"
                ]
            },
            "api": {
                "tech_stack": {
                    "framework": "FastAPI (Python)",
                    "database": "PostgreSQL + SQLAlchemy",
                    "auth": "JWT",
                    "docs": "OpenAPI (Swagger)"
                },
                "folder_structure": [
                    "app/api/v1/endpoints/",
                    "app/core/",
                    "app/models/",
                    "app/schemas/",
                    "app/services/",
                    "tests/",
                    "Dockerfile",
                    "requirements.txt"
                ]
            },
            "cli": {
                "tech_stack": {
                    "language": "Python",
                    "cli_framework": "Click / Typer",
                    "packaging": "pip / PyInstaller"
                },
                "folder_structure": [
                    "src/commands/",
                    "src/utils/",
                    "src/config/",
                    "tests/",
                    "setup.py",
                    "README.md"
                ]
            }
        }

        template = templates.get(app_type, templates["web"])

        plan = {
            "meta": {
                "generated_at": ts,
                "app_type": app_type,
                "description": description,
                "generated_by": "AI ORCHESTRATOR v5.0"
            },
            "tech_stack": template["tech_stack"],
            "folder_structure": template["folder_structure"],
            "development_phases": [
                {
                    "phase": 1,
                    "name": "Setup & Foundation",
                    "tasks": [
                        "Inisialisasi project dan dependencies",
                        "Setup database dan model",
                        "Konfigurasi environment (.env)"
                    ]
                },
                {
                    "phase": 2,
                    "name": "Core Features",
                    "tasks": [
                        "Implementasi fitur utama berdasarkan deskripsi",
                        "API endpoints / route utama",
                        "Business logic layer"
                    ]
                },
                {
                    "phase": 3,
                    "name": "UI/UX",
                    "tasks": [
                        "Design sistem (warna, typography)",
                        "Komponen UI",
                        "Halaman utama dan navigasi"
                    ]
                },
                {
                    "phase": 4,
                    "name": "Integration & Testing",
                    "tasks": [
                        "Integrasi frontend-backend",
                        "Unit tests",
                        "Error handling"
                    ]
                },
                {
                    "phase": 5,
                    "name": "Deploy",
                    "tasks": [
                        "Dockerization",
                        "CI/CD setup",
                        "Dokumentasi README"
                    ]
                }
            ]
        }

        # Simpan ke file jika diminta
        if output_path:
            os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(plan, f, indent=2, ensure_ascii=False)

        # Format output
        output_lines = [
            f"📋 APP PLAN: {description}",
            f"Type: {app_type.upper()} | Generated: {ts}",
            "",
            "🛠️ TECH STACK:",
        ]
        for k, v in template["tech_stack"].items():
            output_lines.append(f"  {k}: {v}")

        output_lines.append("\n📁 FOLDER STRUCTURE:")
        for folder in template["folder_structure"]:
            output_lines.append(f"  {folder}")

        output_lines.append("\n🚀 DEVELOPMENT PHASES:")
        for phase in plan["development_phases"]:
            output_lines.append(f"\n  Phase {phase['phase']}: {phase['name']}")
            for task in phase["tasks"]:
                output_lines.append(f"    - {task}")

        if output_path:
            output_lines.append(f"\n💾 Plan disimpan ke: {output_path}")

        return "\n".join(output_lines)

    except Exception as e:
        return f"Error generating app plan: {e}"
